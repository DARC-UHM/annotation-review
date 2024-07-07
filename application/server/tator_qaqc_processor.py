import datetime
import math
import requests
import sys
import tator

from flask import session
from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .tator_localization_processor import TatorLocalizationProcessor

TERM_RED = '\033[1;31;48m'
TERM_NORMAL = '\033[1;37;0m'


class TatorQaqcProcessor(TatorLocalizationProcessor):
    """
    Fetches annotation information from the Tator given a project id, section id, and list of deployments.
    Filters and formats the annotations for the various QA/QC checks.
    """
    def __init__(
        self,
        project_id: int,
        section_id: int,
        api: tator.api,
        deployment_list: list,
        darc_review_url: str = None,
    ):
        super().__init__(
            project_id=project_id,
            section_id=section_id,
            api=api,
            deployment_list=deployment_list,
            darc_review_url=darc_review_url,
        )

    def fetch_start_times(self):
        for deployment in self.deployments:
            print(f'Fetching media start times for deployment "{deployment}"...', end='')
            sys.stdout.flush()
            if 'media_timestamps' not in session.keys():
                session['media_timestamps'] = {}
            res = requests.get(
                url=f'https://cloud.tator.io/rest/Medias/{self.project_id}?section={self.section_id}&attribute_contains=%24name%3A%3A{deployment}',
                headers={
                    'Content-Type': 'application/json',
                    'Authorization': f'Token {session["tator_token"]}',
                }
            )
            for media in res.json():
                if media['attributes']['Arrival'] and media['attributes']['Arrival'].strip() != '':
                    video_start_timestamp = datetime.fromisoformat(media['attributes']['Start Time'])
                    if 'not observed' in media['attributes']['Arrival']:
                        arrival_frame = 0
                    else:
                        try:
                            arrival_frame = int(media['attributes']['Arrival'].strip().split(' ')[0])
                        except ValueError:
                            print(f'\n{TERM_RED}Error:{TERM_NORMAL} Could not parse Arrival value for {media["name"]}')
                            print(f'Arrival value: "{media["attributes"]["Arrival"]}"')
                            raise ValueError
                    self.bottom_times[deployment] = (video_start_timestamp + datetime.timedelta(seconds=arrival_frame / 30)).strftime('%Y-%m-%d %H:%M:%SZ')
                if media['id'] not in session['media_timestamps'].keys():
                    if 'Start Time' in media['attributes'].keys():
                        session['media_timestamps'][media['id']] = media['attributes']['Start Time']
                        session.modified = True
                    else:
                        print(f'{TERM_RED}Warning:{TERM_NORMAL} No start time found for media {media["id"]}')
            print('fetched!')

    def check_names_accepted(self):
        """
        Finds records with a scientific name or tentative ID that is not accepted in WoRMS
        """
        print('Checking for accepted names...')
        sys.stdout.flush()
        checked = {}
        for localization in self.localizations:
            flag_record = False
            scientific_name = localization['attributes']['Scientific Name']
            tentative_id = localization['attributes']['Tentative ID']
            if scientific_name not in checked.keys():
                if scientific_name in self.phylogeny.keys():
                    checked[scientific_name] = True
                else:
                    if self.fetch_worms_phylogeny(scientific_name):
                        checked[scientific_name] = True
                    else:
                        localization['problems'] = 'Scientific Name'
                        checked[scientific_name] = False
                        flag_record = True
            elif not checked[scientific_name]:
                localization['problems'] = 'Scientific Name'
                flag_record = True
            if tentative_id:
                if tentative_id not in checked.keys():
                    if tentative_id in self.phylogeny.keys():
                        checked[tentative_id] = True
                    else:
                        if self.fetch_worms_phylogeny(tentative_id):
                            checked[tentative_id] = True
                        else:
                            localization['problems'] = 'Tentative ID'
                            checked[tentative_id] = False
                            flag_record = True
                elif not checked[tentative_id]:
                    localization['problems'] = 'Tentative ID' if 'problems' not in localization.keys() else 'Scientific Name, Tentative ID'
                    flag_record = True
            if flag_record:
                self.records_of_interest.append(localization)
        print(f'Found {len(self.records_of_interest)} localizations with unaccepted names!')
        self.save_phylogeny()
        self.process_records(no_match_records={key for key in checked.keys() if not checked[key]})

    def check_missing_qualifier(self):
        """
        Finds records that are classified higher than species but don't have a qualifier set (usually '--'). This check
        need to call process_records first to populate phylogeny.
        """
        self.records_of_interest = self.localizations
        self.process_records()
        actual_final_records = []
        for record in self.final_records:
            if not record['species'] and record['qualifier'] == '--' or not record['qualifier']:
                record['problems'] = 'Scientific Name, Qualifier'
                actual_final_records.append(record)
        self.final_records = actual_final_records

    def check_stet_reason(self):
        """
        Finds records that have a qualifier of 'stet' but no reason set.
        """
        for localization in self.localizations:
            if localization['attributes']['Qualifier'] == 'stet.' and (
                    localization['attributes']['Reason'] == '--' or not localization['attributes']['Reason']):
                localization['problems'] = 'Qualifier, Reason'
                self.records_of_interest.append(localization)
        self.process_records()

    def check_attracted_not_attracted(self, attracted_dict: dict):
        """
        Finds all records that are marked as "attracted" but are saved as "not attracted" in the attracted_dict, and
        vice versa. Also flags all records with taxa that are marked as "attracted/not attracted" in the attracted_dict.
        """
        for localization in self.localizations:
            scientific_name = localization['attributes']['Scientific Name']
            if scientific_name not in attracted_dict.keys() or attracted_dict[scientific_name] == 2:
                localization['problems'] = 'Scientific Name, Attracted'
                self.records_of_interest.append(localization)
            elif localization['attributes']['Attracted'] == 'Attracted' and attracted_dict[scientific_name] == 0:
                localization['problems'] = 'Scientific Name, Attracted'
                self.records_of_interest.append(localization)
            elif localization['attributes']['Attracted'] == 'Not Attracted' and attracted_dict[scientific_name] == 1:
                localization['problems'] = 'Scientific Name, Attracted'
                self.records_of_interest.append(localization)
        self.process_records()

    def check_same_name_qualifier(self):
        """
        Finds records that have the same scientific name/tentative ID combo but a different qualifier.
        """
        scientific_name_qualifiers = {}
        problem_scientific_names = set()
        for localization in self.localizations:
            scientific_name = f'{localization["attributes"]["Scientific Name"]}{" (" + localization["attributes"]["Tentative ID"] + "?)" if localization["attributes"]["Tentative ID"] else ""}'
            if scientific_name not in scientific_name_qualifiers.keys():
                scientific_name_qualifiers[scientific_name] = localization['attributes']['Qualifier']
            else:
                if scientific_name_qualifiers[scientific_name] != localization['attributes']['Qualifier']:
                    problem_scientific_names.add(scientific_name)
        for localization in self.localizations:
            scientific_name = f'{localization["attributes"]["Scientific Name"]}{" (" + localization["attributes"]["Tentative ID"] + "?)" if localization["attributes"]["Tentative ID"] else ""}'
            if scientific_name in problem_scientific_names:
                localization['problems'] = 'Scientific Name, Qualifier'
                self.records_of_interest.append(localization)
        self.process_records()

    def check_non_target_not_attracted(self):
        """
        Finds records that are marked as "non-target" but are marked as "attracted".
        """
        for localization in self.localizations:
            attracted = localization['attributes']['Attracted']
            reason = localization['attributes']['Reason']
            if 'Non-target' in reason and attracted != 'Not Attracted':
                localization['problems'] = 'Attracted, Reason'
                self.records_of_interest.append(localization)
        self.process_records()

    def get_all_tentative_ids(self):
        """
        Finds every record with a tentative ID. Also checks whether or not the tentative ID is in the same
        phylogenetic group as the scientific name.
        """
        no_match_records = set()
        for localization in self.localizations:
            tentative_id = localization['attributes']['Tentative ID']
            if tentative_id and tentative_id not in ['--', '-', '']:
                localization['problems'] = 'Tentative ID'
                self.records_of_interest.append(localization)
        self.process_records()  # process first to make sure phylogeny is populated
        for localization in self.final_records:
            phylogeny_match = False
            if localization['tentative_id'] not in self.phylogeny.keys():
                if localization['tentative_id'] not in no_match_records:
                    if not self.fetch_worms_phylogeny(localization['tentative_id']):
                        no_match_records.add(localization['tentative_id'])
                        localization['problems'] += ' phylogeny no match'
                        continue
                else:
                    localization['problems'] += ' phylogeny no match'
                    continue
            for value in self.phylogeny[localization['tentative_id']].values():
                if value == localization['scientific_name']:
                    phylogeny_match = True
                    break
            if not phylogeny_match:
                localization['problems'] += ' phylogeny no match'
        self.save_phylogeny()

    def get_all_notes_and_remarks(self):
        """
        Finds every record with a note or remark.
        """
        for localization in self.localizations:
            notes = localization['attributes']['Notes']
            id_remarks = localization['attributes']['IdentificationRemarks']
            has_note = notes and notes not in ['--', '-', '']
            has_remark = id_remarks and id_remarks not in ['--', '-', '']
            if has_note and has_remark:
                localization['problems'] = 'Notes, ID Remarks'
                self.records_of_interest.append(localization)
            elif has_note:
                localization['problems'] = 'Notes'
                self.records_of_interest.append(localization)
            elif has_remark:
                localization['problems'] = 'ID Remarks'
                self.records_of_interest.append(localization)
        self.process_records()

    def get_re_examined(self):
        """
        Finds all records that have a reason of "to be re-examined"
        """
        for localization in self.localizations:
            if localization['attributes']['Reason'] == 'To be re-examined':
                self.records_of_interest.append(localization)
        self.process_records()

    def get_unique_taxa(self):
        """
        Finds every unique scientific name/tentative ID combo and box/dot info.
        """
        self.fetch_start_times()
        self.records_of_interest = self.localizations
        self.process_records(get_timestamp=True)
        unique_taxa = {}
        for record in self.final_records:
            scientific_name = record['scientific_name']
            tentative_id = record['tentative_id']
            key = f'{scientific_name}:{tentative_id}'
            if key not in unique_taxa.keys():
                # add new unique taxa to dict
                unique_taxa[key] = {
                    'scientific_name': scientific_name,
                    'tentative_id': tentative_id,
                    'box_count': 0,
                    'dot_count': 0,
                    'first_box': '',
                    'first_dot': '',
                }
            for localization in record['all_localizations']:
                # increment box/dot counts, set first box/dot and TOFA
                if localization['type'] == 48:
                    unique_taxa[key]['box_count'] += 1
                    first_box = unique_taxa[key]['first_box']
                    if not first_box or datetime.strptime(record['timestamp'], '%Y-%m-%d %H:%M:%SZ') < datetime.strptime(first_box, '%Y-%m-%d %H:%M:%SZ'):
                        unique_taxa[key]['first_box'] = record['timestamp']
                        unique_taxa[key]['first_box_url'] = f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}&selected_entity={localization["elemental_id"]}'
                elif localization['type'] == 49:
                    unique_taxa[key]['dot_count'] += 1
                    first_dot = unique_taxa[key]['first_dot']
                    observed_timestamp = datetime.strptime(record['timestamp'], '%Y-%m-%d %H:%M:%SZ')
                    if not first_dot or observed_timestamp < datetime.strptime(first_dot, '%Y-%m-%d %H:%M:%SZ'):
                        unique_taxa[key]['first_dot'] = record['timestamp']
                        unique_taxa[key]['first_dot_url'] = f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}&selected_entity={localization["elemental_id"]}'
        self.final_records = unique_taxa

    def get_max_n(self):
        """
        Finds the highest dot count for each unique scientific name/tentative ID combo per deployment.
        """
        self.records_of_interest = self.localizations
        self.process_records(get_ctd=True)
        deployment_taxa = {}
        unique_taxa = {}
        for record in self.final_records:
            scientific_tentative = f'{record["scientific_name"]}{" (" + record["tentative_id"] + "?)" if record["tentative_id"] else ""}'
            if record['count'] < 1 or record['attracted'] == 'Not Attracted':
                continue
            if scientific_tentative not in unique_taxa.keys():
                unique_taxa[scientific_tentative] = {
                    'scientific_tentative': scientific_tentative,
                    'phylum': record.get('phylum'),
                    'class': record.get('class'),
                    'order': record.get('order'),
                    'family': record.get('family'),
                    'genus': record.get('genus'),
                    'species': record.get('species'),
                }
            if record['video_sequence_name'] not in deployment_taxa.keys():
                deployment_taxa[record['video_sequence_name']] = {
                    'depth_m': record['depth_m'],
                    'max_n_dict': {},
                }
            if scientific_tentative not in deployment_taxa[record['video_sequence_name']]['max_n_dict'].keys():
                # add new unique taxa to dict
                deployment_taxa[record['video_sequence_name']]['max_n_dict'][scientific_tentative] = {
                    'max_n': record['count'],
                    'max_n_url': f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}',
                }
            else:
                # check for new max N
                if record['count'] > deployment_taxa[record['video_sequence_name']]['max_n_dict'][scientific_tentative]['max_n']:
                    deployment_taxa[record['video_sequence_name']]['max_n_dict'][scientific_tentative]['max_n'] = record['count']
                    deployment_taxa[record['video_sequence_name']]['max_n_dict'][scientific_tentative]['max_n_url'] = f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}'
        # convert unique taxa to list for sorting
        unique_taxa_list = list(unique_taxa.values())
        unique_taxa_list.sort(key=lambda x: (
            x['phylum'] if x.get('phylum') else '',
            x['class'] if x.get('class') else '',
            x['order'] if x.get('order') else '',
            x['family'] if x.get('family') else '',
            x['genus'] if x.get('genus') else '',
            x['species'] if x.get('species') else '',
        ))
        self.final_records = {'deployments': deployment_taxa, 'unique_taxa': [taxa['scientific_tentative'] for taxa in unique_taxa_list]}

    def get_tofa(self):
        """
        Finds the time of first arrival for each unique scientific name/tentative ID combo per deployment. Also shows
        species accumulation curve. Ignores non-attracted taxa.
        """
        self.fetch_start_times()
        self.records_of_interest = self.localizations
        self.process_records(get_timestamp=True, get_ctd=True)
        deployment_taxa = {}
        unique_taxa = {}
        unique_taxa_first_seen = {}
        bottom_time = None
        latest_timestamp = datetime.fromtimestamp(0)  # to find the duration of the deployment
        for record in self.final_records:
            scientific_tentative = f'{record["scientific_name"]}{" (" + record["tentative_id"] + "?)" if record["tentative_id"] else ""}'
            observed_timestamp = datetime.strptime(record['timestamp'], '%Y-%m-%d %H:%M:%SZ')
            bottom_time = datetime.strptime(self.bottom_times[record['video_sequence_name']], '%Y-%m-%d %H:%M:%SZ')
            if record['count'] < 1 or record['attracted'] == 'Not Attracted':
                continue
            if observed_timestamp > latest_timestamp:
                latest_timestamp = observed_timestamp
            if scientific_tentative not in unique_taxa_first_seen.keys():
                unique_taxa_first_seen[scientific_tentative] = observed_timestamp
            else:
                if observed_timestamp < unique_taxa_first_seen[scientific_tentative]:
                    unique_taxa_first_seen[scientific_tentative] = observed_timestamp
            if scientific_tentative not in unique_taxa.keys():
                unique_taxa[scientific_tentative] = {
                    'scientific_tentative': scientific_tentative,
                    'phylum': record.get('phylum'),
                    'class': record.get('class'),
                    'order': record.get('order'),
                    'family': record.get('family'),
                    'genus': record.get('genus'),
                    'species': record.get('species'),
                }
            if record['video_sequence_name'] not in deployment_taxa.keys():
                deployment_taxa[record['video_sequence_name']] = {
                    'depth_m': record['depth_m'],
                    'tofa_dict': {},
                }
            if scientific_tentative not in deployment_taxa[record['video_sequence_name']]['tofa_dict'].keys():
                # add new unique taxa to dict
                deployment_taxa[record['video_sequence_name']]['tofa_dict'][scientific_tentative] = {
                    'tofa': str(observed_timestamp - bottom_time) if observed_timestamp > bottom_time else '00:00:00',
                    'tofa_url': f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}',
                }
            else:
                # check for new tofa
                if str(observed_timestamp - bottom_time) < deployment_taxa[record['video_sequence_name']]['tofa_dict'][scientific_tentative]['tofa']:
                    deployment_taxa[record['video_sequence_name']]['tofa_dict'][scientific_tentative]['tofa'] = str(observed_timestamp - bottom_time) if observed_timestamp > bottom_time else '00:00:00'
                    deployment_taxa[record['video_sequence_name']]['tofa_dict'][scientific_tentative]['tofa_url'] = f'https://cloud.tator.io/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}'
        # convert unique taxa to list for sorting
        unique_taxa_list = list(unique_taxa.values())
        unique_taxa_list.sort(key=lambda x: (
            x['phylum'] if x.get('phylum') else '',
            x['class'] if x.get('class') else '',
            x['order'] if x.get('order') else '',
            x['family'] if x.get('family') else '',
            x['genus'] if x.get('genus') else '',
            x['species'] if x.get('species') else '',
        ))
        # rounding up to nearest hour
        deployment_time = datetime.timedelta(hours=math.ceil((latest_timestamp - bottom_time).total_seconds() / 3600))
        accumulation_data = []  # just a list of the number of unique taxa seen at each hour
        for hour in range(1, deployment_time.seconds // 3600 + 1):
            accumulation_data.append(len([taxa for taxa in unique_taxa_first_seen.values() if taxa < bottom_time + datetime.timedelta(hours=hour)]))
        self.final_records = {
            'deployments': deployment_taxa,
            'unique_taxa': [taxa['scientific_tentative'] for taxa in unique_taxa_list],
            'deployment_time': deployment_time.seconds // 3600,
            'accumulation_data': accumulation_data,
        }

    def get_summary(self):
        """
        Returns a summary of the final records.
        """
        self.fetch_start_times()
        self.records_of_interest = [localization for localization in self.localizations if localization['type'] != 48]
        self.process_records(get_timestamp=True, get_ctd=True, get_substrates=True)

    def download_image_guide(self, app) -> Presentation:
        """
        Finds all records marked as "good" images, saves them to a ppt.
        """
        for localization in self.localizations:
            if localization['attributes'].get('Good Image'):
                self.records_of_interest.append(localization)
        self.process_records()
        pres = Presentation()
        image_slide_layout = pres.slide_layouts[6]

        i = 0
        while i < len(self.final_records):
            slide = pres.slides.add_slide(image_slide_layout)
            current_phylum = self.final_records[i].get('phylum')
            if current_phylum is None:
                current_phylum = 'UNKNOWN PHYLUM'
            phylum_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            phylum_text_frame = phylum_text_box.text_frame
            phylum_paragraph = phylum_text_frame.paragraphs[0]
            phylum_paragraph.alignment = PP_ALIGN.CENTER
            phylum_run = phylum_paragraph.add_run()
            phylum_run.text = ' '.join(list(current_phylum.upper()))
            phylum_font = phylum_run.font
            phylum_font.name = 'Arial'
            phylum_font.size = Pt(32)
            phylum_font.color.rgb = RGBColor(0, 0, 0)
            for j in range(4):
                # add four images to slide
                localization = self.final_records[i]
                if localization['phylum'] != current_phylum and current_phylum != 'UNKNOWN PHYLUM':
                    break
                localization_id = localization['all_localizations'][0]['id']
                response = requests.get(f'{app.config.get("LOCAL_APP_URL")}/tator-localization/{localization_id}?token={session["tator_token"]}')
                if response.status_code != 200:
                    print(f'Error fetching image for record {localization["observation_uuid"]}')
                    continue
                image_data = BytesIO(response.content)
                top = Inches(1.5 if j < 2 else 4)
                left = Inches(1 if j % 2 == 0 else 5)
                picture = slide.shapes.add_picture(image_data, left, top, height=Inches(2.5))
                line = picture.line
                line.color.rgb = RGBColor(0, 0, 0)
                line.width = Pt(1.5)
                # add text box
                width = Inches(2)
                height = Inches(1)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                paragraph = text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = f'{localization["scientific_name"]}{" (" + localization["tentative_id"] + "?)" if localization.get("tentative_id") else ""}'
                font = run.font
                font.name = 'Arial'
                font.size = Pt(18)
                font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                font.italic = True
                if localization['attracted'] == 'Not Attracted':
                    text_frame.add_paragraph()
                    paragraph = text_frame.paragraphs[1]
                    run_2 = paragraph.add_run()
                    run_2.text = 'NOT ATTRACTED'
                    font = run_2.font
                    font.name = 'Arial'
                    font.size = Pt(18)
                    font.color.rgb = RGBColor(0xff, 0x0, 0x0)
                    font.italic = False
                i += 1
                if i >= len(self.final_records):
                    break
        return pres
