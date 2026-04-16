from io import BytesIO

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from application.tator.tator_rest_client import TatorRestClient
from application.tator.tator_type import TatorLocalizationType


class ImageGuidePresentation:
    """
    Builds a PowerPoint image guide from a list of processed Tator localization records.
    Six images per slide (3x2 grid), each with a black header bar, grouped by phylum.
    """

    IMAGE_ASPECT_RATIO = 16 / 9
    IMAGE_WIDTH = Inches(3.0)
    IMAGE_HEIGHT = Inches(3.0 / IMAGE_ASPECT_RATIO)
    IMAGE_HEADER_HEIGHT = Inches(0.45)
    ROW_TOPS = [Inches(1.5), Inches(4.25)]  # header top for each row
    BORDER_WIDTH = Pt(1.5)

    def __init__(self, tator_client: TatorRestClient):
        self.tator_client = tator_client

    def build(self, records: list[dict]):
        pres = Presentation()
        image_slide_layout = pres.slide_layouts[6]

        i = 0
        while i < len(records):
            slide = pres.slides.add_slide(image_slide_layout)
            current_phylum = records[i].get('phylum') or 'UNKNOWN PHYLUM'
            self._add_phylum_header(slide, current_phylum)
            for j in range(6):
                if i >= len(records):
                    break
                localization = records[i]
                if localization.get('phylum') != current_phylum and current_phylum != 'UNKNOWN PHYLUM':
                    break
                print(f'Processing image {i + 1}/{len(records)}')
                try:
                    image_data = self._fetch_normalized_image(localization)
                except Exception as e:
                    print(f'Error fetching image for localization {localization["id"]}: {e}')
                    i += 1
                    continue
                header_top = self.ROW_TOPS[j // 3]
                image_top = header_top + self.IMAGE_HEADER_HEIGHT
                left = Inches(0.5 + (j % 3) * 3.0)
                self._add_image_header(slide, localization, left, header_top)
                self._add_image(slide, image_data, left, image_top)
                if (TatorLocalizationType.is_dropcam(localization.get('type')) and
                        (localization.get('attracted') is None or localization['attracted'] == 'Not Attracted')):
                    self._add_not_attracted_overlay(slide, left, image_top)
                i += 1
        return pres

    def _add_phylum_header(self, slide, phylum: str):
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = ' '.join(list(phylum.upper()))
        font = run.font
        font.name = 'Arial'
        font.size = Pt(32)
        font.color.rgb = RGBColor(0, 0, 0)

    def _add_image(self, slide, image_data: BytesIO, left, top):
        picture = slide.shapes.add_picture(image_data, left, top, width=self.IMAGE_WIDTH, height=self.IMAGE_HEIGHT)
        picture.line.color.rgb = RGBColor(0, 0, 0)
        picture.line.width = self.BORDER_WIDTH

    def _add_image_header(self, slide, localization: dict, left, top):
        # Black background rectangle
        rect = slide.shapes.add_shape(
            autoshape_type_id=MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left=left - self.BORDER_WIDTH,
            top=top,
            width=self.IMAGE_WIDTH + 2 * self.BORDER_WIDTH,
            height=self.IMAGE_HEADER_HEIGHT,
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
        rect.line.fill.background()

        # Text box over the rectangle
        text_box = slide.shapes.add_textbox(left, top, self.IMAGE_WIDTH, self.IMAGE_HEADER_HEIGHT)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]

        paragraph.alignment = PP_ALIGN.CENTER
        tentative_id = localization.get('tentative_id')
        morphospecies = localization.get('morphospecies')
        self._make_run(paragraph, localization['scientific_name'], italic=bool(localization.get('genus')))
        if localization.get('genus') and not localization.get('species'):
            self._make_run(paragraph, ' sp.', italic=False)
        if tentative_id or morphospecies:
            extra_id = tentative_id or morphospecies
            self._make_run(paragraph, ' (', italic=False)
            self._make_run(paragraph, extra_id, italic=bool(localization.get('family')))
            if localization.get('family') and not morphospecies:
                self._make_run(paragraph, ' sp.', italic=False)
            if tentative_id:
                self._make_run(paragraph, '?)', italic=False)
            else:
                self._make_run(paragraph, ')', italic=False)

    def _add_not_attracted_overlay(self, slide, left, image_top):
        overlay_height = Inches(0.35)
        text_box = slide.shapes.add_textbox(left, image_top + self.IMAGE_HEIGHT - overlay_height, self.IMAGE_WIDTH, overlay_height)
        text_frame = text_box.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = 'NOT ATTRACTED'
        font = run.font
        font.name = 'Arial'
        font.size = Pt(14)
        font.color.rgb = RGBColor(0xff, 0x0, 0x0)
        font.bold = True

    @staticmethod
    def _make_run(paragraph, text, italic):
        run = paragraph.add_run()
        run.text = text
        run.font.name = 'Arial'
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        run.font.italic = italic
        return run

    def _fetch_normalized_image(self, localization: dict) -> BytesIO:
        """Fetches full frame from Tator, crops to localization bounds, and expands to 16:9 aspect ratio."""
        frame_bytes = self.tator_client.get_frame(localization['media_id'], frame=localization['frame'])
        img = Image.open(BytesIO(frame_bytes))
        img_width, img_height = img.width, img.height

        box_localization = localization['all_localizations'][0]  # there should only ever be one localization (the box)
        x, y = box_localization['points']
        w, h = box_localization['dimensions']
        left = int(x * img_width)
        upper = int(y * img_height)
        right = int((x + w) * img_width)
        lower = int((y + h) * img_height)

        # expand to 16:9 aspect ratio
        aspect_ratio = (right - left) / (lower - upper)
        if aspect_ratio < self.IMAGE_ASPECT_RATIO:
            delta = ((lower - upper) * self.IMAGE_ASPECT_RATIO - (right - left)) / 2
            left -= delta
            right += delta
        elif aspect_ratio > self.IMAGE_ASPECT_RATIO:
            delta = ((right - left) / self.IMAGE_ASPECT_RATIO - (lower - upper)) / 2
            upper -= delta
            lower += delta

        left, upper, right, lower = int(left), int(upper), int(right), int(lower)
        if left < 0:
            right += abs(left)
            left = 0
        if upper < 0:
            lower += abs(upper)
            upper = 0
        if right > img_width:
            left -= right - img_width
            right = img_width
        if lower > img_height:
            upper -= lower - img_height
            lower = img_height

        img = img.crop((left, upper, right, lower))
        output = BytesIO()
        img.save(output, format='JPEG')
        output.seek(0)
        return output
