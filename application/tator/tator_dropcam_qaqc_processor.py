import datetime
import math
import sys
from io import BytesIO

import requests
import tator
from flask import session
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from application.tator.tator_base_qaqc_processor import TatorBaseQaqcProcessor
from application.util.constants import TERM_NORMAL, TERM_RED
from application.tator.tator_type import TatorLocalizationType


class TatorDropcamQaqcProcessor(TatorBaseQaqcProcessor):
    def __init__(
            self,
            project_id: int,
            section_ids: list[str],
            api: tator.api,
            tator_url: str,
            darc_review_url: str = None,
            transect_media_ids: list[int] = None,
    ):
        super().__init__(
            project_id=project_id,
            section_ids=section_ids,
            api=api,
            darc_review_url=darc_review_url,
            tator_url=tator_url,
            transect_media_ids=transect_media_ids,
        )

    def check_attracted_not_attracted(self, attracted_dict: dict):
        """
        Finds all records that are marked as "attracted" but are saved as "not attracted" in the attracted_dict, and
        vice versa. Also flags all records with taxa that are marked as "attracted/not attracted" in the attracted_dict.
        """
        for section in self.sections:
            records_of_interest = []
            for localization in section.localizations:
                scientific_name = localization['attributes'].get('Scientific Name')
                if scientific_name not in attracted_dict.keys() or attracted_dict[scientific_name] == 2:
                    localization['problems'] = 'Scientific Name, Attracted'
                    records_of_interest.append(localization)
                elif localization['attributes'].get('Attracted') == 'Attracted' and attracted_dict[scientific_name] == 0:
                    localization['problems'] = 'Scientific Name, Attracted'
                    records_of_interest.append(localization)
                elif localization['attributes'].get('Attracted') == 'Not Attracted' and attracted_dict[scientific_name] == 1:
                    localization['problems'] = 'Scientific Name, Attracted'
                    records_of_interest.append(localization)
            section.localizations = records_of_interest
        self.process_records()

    def check_same_name_qualifier(self):
        """
        Finds records that have the same scientific name/tentative ID combo but a different qualifier.
        """
        scientific_name_qualifiers = {}
        problem_scientific_names = set()
        for section in self.sections:
            # first pass: build dict of scientific name/tentative ID combos and their qualifiers, add to problem set if mismatch
            for localization in section.localizations:
                scientific_name = f'{localization["attributes"].get("Scientific Name")}{" (" + localization["attributes"]["Tentative ID"] + "?)" if localization["attributes"].get("Tentative ID") else ""}'
                if scientific_name not in scientific_name_qualifiers.keys():
                    scientific_name_qualifiers[scientific_name] = localization['attributes'].get('Qualifier')
                else:
                    if scientific_name_qualifiers[scientific_name] != localization['attributes'].get('Qualifier'):
                        problem_scientific_names.add(scientific_name)
        for section in self.sections:
            # second pass: add records with problem scientific names to records of interest
            records_of_interest = []
            for localization in section.localizations:
                scientific_name = f'{localization["attributes"].get("Scientific Name")}{" (" + localization["attributes"]["Tentative ID"] + "?)" if localization["attributes"].get("Tentative ID") else ""}'
                if scientific_name in problem_scientific_names:
                    localization['problems'] = 'Scientific Name, Qualifier'
                    records_of_interest.append(localization)
            section.localizations = records_of_interest
        self.process_records()

    def check_non_target_not_attracted(self):
        """
        Finds records that are marked as "non-target" but are marked as "attracted".
        """
        for section in self.sections:
            records_of_interest = []
            for localization in section.localizations:
                attracted = localization['attributes'].get('Attracted')
                reason = localization['attributes'].get('Reason')
                if 'Non-target' in reason and attracted != 'Not Attracted':
                    localization['problems'] = 'Attracted, Reason'
                    records_of_interest.append(localization)
            section.localizations = records_of_interest
        self.process_records()

    def check_exists_in_image_references(self, image_refs: dict):
        """
        Finds records that do not exist in the image references db (combo scientific name, tentative ID,
        and morphospecies). Also flags records with both tentative ID and morphospecies set.
        """
        for section in self.sections:
            records_of_interest = []
            for localization in section.localizations:
                image_ref_key = localization['attributes'].get('Scientific Name')
                tentative_id = localization['attributes'].get('Tentative ID')
                morphospecies = localization['attributes'].get('Morphospecies')
                if tentative_id and morphospecies:
                    localization['problems'] = 'Tentative ID, Morphospecies'
                    records_of_interest.append(localization)
                    continue
                if tentative_id and tentative_id != '':
                    image_ref_key += f'~tid={tentative_id}'
                if morphospecies and morphospecies != '':
                    image_ref_key += f'~m={morphospecies}'
                if image_ref_key not in image_refs:
                    records_of_interest.append(localization)
            section.localizations = records_of_interest
        self.process_records()

    def get_unique_taxa(self):
        self.fetch_start_times()
        self.process_records(get_timestamp=True)
        unique_taxa = {}
        for record in self.final_records:
            scientific_name = record.get('scientific_name')
            tentative_id = record.get('tentative_id', '')
            morphospecies = record.get('morphospecies', '')
            key = f'{scientific_name}:{tentative_id}:{morphospecies}'
            if key not in unique_taxa.keys():
                # add new unique taxa to dict
                unique_taxa[key] = {
                    'scientific_name': scientific_name,
                    'tentative_id': tentative_id,
                    'morphospecies': morphospecies,
                    'box_count': 0,
                    'dot_count': 0,
                    'first_box': '',
                    'first_dot': '',
                }
            for localization in record['all_localizations']:
                # increment box/dot counts, set first box/dot and TOFA
                if TatorLocalizationType.is_box(localization['type']):
                    unique_taxa[key]['box_count'] += 1
                    if not record.get('timestamp'):
                        continue
                    first_box = unique_taxa[key]['first_box']
                    observed_timestamp = datetime.datetime.strptime(record['timestamp'], self.BOTTOM_TIME_FORMAT)
                    if not first_box or observed_timestamp < datetime.datetime.strptime(first_box, self.BOTTOM_TIME_FORMAT):
                        unique_taxa[key]['first_box'] = record['timestamp']
                        unique_taxa[key]['first_box_url'] = f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}&selected_entity={localization["elemental_id"]}'
                elif TatorLocalizationType.is_dot(localization['type']):
                    unique_taxa[key]['dot_count'] += 1
                    if not record.get('timestamp'):
                        continue
                    first_dot = unique_taxa[key]['first_dot']
                    observed_timestamp = datetime.datetime.strptime(record['timestamp'], self.BOTTOM_TIME_FORMAT)
                    if not first_dot or observed_timestamp < datetime.datetime.strptime(first_dot, self.BOTTOM_TIME_FORMAT):
                        unique_taxa[key]['first_dot'] = record['timestamp']
                        unique_taxa[key]['first_dot_url'] = f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}&selected_entity={localization["elemental_id"]}'
        self.final_records = unique_taxa

    def get_max_n(self):
        """
        Finds the highest dot count for each unique scientific name, tentative ID, and morphospecies combo per
        deployment. Ignores non-attracted taxa.
        """
        self.process_records(get_ctd=True)
        deployment_taxa = {}
        unique_taxa = {}
        for record in self.final_records:
            scientific_name = record.get('scientific_name')
            tentative_id_suffix = f' ({record["tentative_id"]}?)' if record.get('tentative_id') else ''
            morphospecies_suffix = f' ({record["morphospecies"]})' if record.get('morphospecies') else ''
            unique_name = f'{scientific_name}{tentative_id_suffix}{morphospecies_suffix}'
            if record.get('count', 0) < 1 or record.get('attracted') == 'Not Attracted':
                continue
            if unique_name not in unique_taxa.keys():
                unique_taxa[unique_name] = {
                    'unique_name': unique_name,
                    'phylum': record.get('phylum'),
                    'class': record.get('class'),
                    'order': record.get('order'),
                    'family': record.get('family'),
                    'genus': record.get('genus'),
                    'species': record.get('species'),
                }
            if record['video_sequence_name'] not in deployment_taxa.keys():
                deployment_taxa[record['video_sequence_name']] = {
                    'depth_m': record.get('depth_m'),
                    'max_n_dict': {},
                }
            if unique_name not in deployment_taxa[record['video_sequence_name']]['max_n_dict'].keys():
                # add new unique taxa to dict
                deployment_taxa[record['video_sequence_name']]['max_n_dict'][unique_name] = {
                    'max_n': record['count'],
                    'max_n_url': f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}',
                }
            else:
                # check for new max N
                if record['count'] > deployment_taxa[record['video_sequence_name']]['max_n_dict'][unique_name]['max_n']:
                    deployment_taxa[record['video_sequence_name']]['max_n_dict'][unique_name]['max_n'] = record['count']
                    deployment_taxa[record['video_sequence_name']]['max_n_dict'][unique_name]['max_n_url'] = f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}'
        # convert unique taxa to list for sorting
        unique_taxa_list = list(unique_taxa.values())
        unique_taxa_list.sort(key=lambda x: (
            x['phylum'] if x.get('phylum') else '',
            x['class'] if x.get('class') else '',
            x['order'] if x.get('order') else '',
            x['family'] if x.get('family') else '',
            x['genus'] if x.get('genus') else '',
            x['species'] if x.get('species') else '',
        ))
        self.final_records = {
            'deployments': deployment_taxa,
            'unique_taxa': [taxa['unique_name'] for taxa in unique_taxa_list],
        }

    def get_tofa(self):
        """
        Finds the time of first arrival for each unique scientific name, tentative ID, and morphospecies combo per
        deployment. Also shows species accumulation curve. Ignores non-attracted taxa.
        """
        self.fetch_start_times()
        self.process_records(get_timestamp=True, get_ctd=True)
        deployment_taxa = {}
        unique_taxa = {}
        unique_taxa_first_seen = {}
        section_id_indices = {section.section_id: index for index, section in enumerate(self.sections)}
        bottom_time = None
        latest_timestamp = datetime.datetime.fromtimestamp(0)  # to find the duration of the deployment
        for record in self.final_records:
            scientific_name = record.get('scientific_name')
            tentative_id_suffix = f' ({record["tentative_id"]}?)' if record.get('tentative_id') else ''
            morphospecies_suffix = f' ({record["morphospecies"]})' if record.get('morphospecies') else ''
            unique_name = f'{scientific_name}{tentative_id_suffix}{morphospecies_suffix}'
            if not record.get('timestamp'):
                continue
            observed_timestamp = datetime.datetime.strptime(record['timestamp'], self.BOTTOM_TIME_FORMAT)
            this_section = self.sections[section_id_indices[record['section_id']]]
            bottom_time = datetime.datetime.strptime(this_section.bottom_time, self.BOTTOM_TIME_FORMAT)
            if record.get('count', 0) < 1 or record.get('attracted') == 'Not Attracted':
                continue
            if observed_timestamp > latest_timestamp:
                latest_timestamp = observed_timestamp
            if unique_name not in unique_taxa_first_seen.keys():
                unique_taxa_first_seen[unique_name] = observed_timestamp
            else:
                if observed_timestamp < unique_taxa_first_seen[unique_name]:
                    unique_taxa_first_seen[unique_name] = observed_timestamp
            if unique_name not in unique_taxa.keys():
                unique_taxa[unique_name] = {
                    'unique_name': unique_name,
                    'phylum': record.get('phylum'),
                    'class': record.get('class'),
                    'order': record.get('order'),
                    'family': record.get('family'),
                    'genus': record.get('genus'),
                    'species': record.get('species'),
                }
            if record['video_sequence_name'] not in deployment_taxa.keys():
                deployment_taxa[record['video_sequence_name']] = {
                    'depth_m': record.get('depth_m'),
                    'tofa_dict': {},
                }
            time_diff = observed_timestamp - bottom_time if observed_timestamp > bottom_time else datetime.timedelta(0)
            if unique_name not in deployment_taxa[record['video_sequence_name']]['tofa_dict'].keys():
                # add new unique taxa to dict
                deployment_taxa[record['video_sequence_name']]['tofa_dict'][unique_name] = {
                    'tofa': str(time_diff),
                    'tofa_seconds': time_diff.total_seconds(),
                    'tofa_url': f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}',
                }
            else:
                # check for new tofa
                if time_diff.total_seconds() < deployment_taxa[record['video_sequence_name']]['tofa_dict'][unique_name]['tofa_seconds']:
                    deployment_taxa[record['video_sequence_name']]['tofa_dict'][unique_name]['tofa'] = str(time_diff)
                    deployment_taxa[record['video_sequence_name']]['tofa_dict'][unique_name]['tofa_seconds'] = time_diff.total_seconds()
                    deployment_taxa[record['video_sequence_name']]['tofa_dict'][unique_name]['tofa_url'] = \
                        f'{self.tator_url}/{self.project_id}/annotation/{record["media_id"]}?frame={record["frame"]}'
        # convert unique taxa to list for sorting
        unique_taxa_list = list(unique_taxa.values())
        unique_taxa_list.sort(key=lambda x: (
            x['phylum'] if x.get('phylum') else '',
            x['class'] if x.get('class') else '',
            x['order'] if x.get('order') else '',
            x['family'] if x.get('family') else '',
            x['genus'] if x.get('genus') else '',
            x['species'] if x.get('species') else '',
        ))
        if len(unique_taxa_list) == 0:
            print(f'{TERM_RED}ERROR: Unable to calculate TOFA. Missing start times?{TERM_NORMAL}')
            self.final_records = {
                'deployments': deployment_taxa,
                'unique_taxa': [],
                'deployment_time': 0,
                'accumulation_data': [],
            }
            return
        # rounding up to nearest hour
        deployment_time = datetime.timedelta(hours=math.ceil((latest_timestamp - bottom_time).total_seconds() / 3600))
        accumulation_data = []  # just a list of the number of unique taxa seen at each hour
        for hour in range(1, deployment_time.seconds // 3600 + 1):
            accumulation_data.append(len([
                taxa for taxa in unique_taxa_first_seen.values() if taxa < bottom_time + datetime.timedelta(hours=hour)
            ]))
        self.final_records = {
            'deployments': deployment_taxa,
            'unique_taxa': [taxa['unique_name'] for taxa in unique_taxa_list],
            'deployment_time': deployment_time.seconds // 3600,
            'accumulation_data': accumulation_data,
        }

    def get_summary(self):
        self.fetch_start_times()
        for section in self.sections:
            section.localizations = [
                localization for localization in section.localizations if not TatorLocalizationType.is_box(localization['type'])
            ]
        self.process_records(get_timestamp=True, get_ctd=True, get_substrates=True)

    def download_image_guide(self, app) -> Presentation:
        for section in self.sections:
            records_of_interest = []
            for localization in section.localizations:
                if localization['attributes'].get('Good Image'):
                    records_of_interest.append(localization)
            section.localizations = records_of_interest
        self.process_records()
        pres = Presentation()
        image_slide_layout = pres.slide_layouts[6]

        i = 0
        while i < len(self.final_records):
            slide = pres.slides.add_slide(image_slide_layout)
            current_phylum = self.final_records[i].get('phylum')
            if current_phylum is None:
                current_phylum = 'UNKNOWN PHYLUM'
            phylum_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            phylum_text_frame = phylum_text_box.text_frame
            phylum_paragraph = phylum_text_frame.paragraphs[0]
            phylum_paragraph.alignment = PP_ALIGN.CENTER
            phylum_run = phylum_paragraph.add_run()
            phylum_run.text = ' '.join(list(current_phylum.upper()))
            phylum_font = phylum_run.font
            phylum_font.name = 'Arial'
            phylum_font.size = Pt(32)
            phylum_font.color.rgb = RGBColor(0, 0, 0)
            for j in range(4):
                # add four images to slide
                localization = self.final_records[i]
                if localization['phylum'] != current_phylum and current_phylum != 'UNKNOWN PHYLUM':
                    break
                localization_id = localization['all_localizations'][0]['id']
                response = requests.get(f'{app.config.get("LOCAL_APP_URL")}/tator/localization-image/{localization_id}?token={session["tator_token"]}')
                if response.status_code != 200:
                    print(f'Error fetching image for record {localization["observation_uuid"]}')
                    continue
                image_data = BytesIO(response.content)
                top = Inches(1.5 if j < 2 else 4)
                left = Inches(1 if j % 2 == 0 else 5)
                picture = slide.shapes.add_picture(image_data, left, top, height=Inches(2.5))
                line = picture.line
                line.color.rgb = RGBColor(0, 0, 0)
                line.width = Pt(1.5)
                # add text box
                width = Inches(2)
                height = Inches(1)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                paragraph = text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = f'{localization["scientific_name"]}{" (" + localization["tentative_id"] + "?)" if localization.get("tentative_id") else ""}'
                font = run.font
                font.name = 'Arial'
                font.size = Pt(18)
                font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                font.italic = True
                if localization['attracted'] == 'Not Attracted':
                    text_frame.add_paragraph()
                    paragraph = text_frame.paragraphs[1]
                    run_2 = paragraph.add_run()
                    run_2.text = 'NOT ATTRACTED'
                    font = run_2.font
                    font.name = 'Arial'
                    font.size = Pt(18)
                    font.color.rgb = RGBColor(0xff, 0x0, 0x0)
                    font.italic = False
                i += 1
                if i >= len(self.final_records):
                    break
        return pres

    def fetch_start_times(self):
        if 'media_timestamps' not in session.keys():
            session['media_timestamps'] = {}
        for section in self.sections:
            print(f'Fetching media start times for deployment "{section.deployment_name}"...', end='')
            sys.stdout.flush()
            for media in self.tator_client.get_medias_for_section(self.project_id, section=section.section_id):
                # get media start times
                if media['id'] not in session['media_timestamps'].keys():
                    if 'Start Time' in media['attributes'].keys():
                        session['media_timestamps'][media['id']] = media['attributes']['Start Time']
                        session.modified = True
                    else:
                        print(f'{TERM_RED}Warning:{TERM_NORMAL} No start time found for media {media["id"]}')
                        continue
                # get deployment bottom time
                media_arrival_attribute = media['attributes'].get('Arrival')
                if media_arrival_attribute and media_arrival_attribute.strip() != '':
                    video_start_timestamp = datetime.datetime.fromisoformat(media['attributes']['Start Time']).astimezone(datetime.timezone.utc)
                    if 'not observed' in media_arrival_attribute.lower():
                        arrival_frame = 0
                    else:
                        try:
                            arrival_frame = int(media_arrival_attribute.strip().split(' ')[0])
                        except ValueError:
                            error_message = (f'Could not parse Arrival value for media "{media["name"]}". '
                                             f'Expected format like "1234" or "not observed" but got "{media["attributes"]["Arrival"]}".')
                            print(f'\n\n{TERM_RED}ERROR: {error_message}{TERM_NORMAL}')
                            raise ValueError(error_message)
                    media_fps = media.get('fps') or 30
                    deployment_bottom_time = video_start_timestamp + datetime.timedelta(seconds=arrival_frame / media_fps)
                    section.bottom_time = deployment_bottom_time.strftime(self.BOTTOM_TIME_FORMAT)
            print('fetched!')
