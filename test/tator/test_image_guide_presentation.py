import time
from io import BytesIO
from unittest.mock import patch

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from application.tator.image_guide_presentation import ImageGuidePresentation
from application.tator.tator_rest_client import TatorRestClient
from application.tator.tator_type import TatorLocalizationType


def make_slide():
    presentation = Presentation()
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def make_jpeg_bytes(width=100, height=100, color='red'):
    buffer = BytesIO()
    Image.new('RGB', (width, height), color=color).save(buffer, format='JPEG')
    buffer.seek(0)
    return buffer


def run_texts(paragraph):
    return [(run.text, run.font.italic) for run in paragraph.runs]


def pictures(slide):
    return [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]


def text_boxes_with_text(slide, text):
    return [
        shape for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text == text
    ]


@pytest.fixture
def image_guide_presentation():
    return ImageGuidePresentation(tator_client=TatorRestClient('https://tator.url', 'super-secret-token'))


class TestImageGuidePresentation:

    def test_build_groups_records_into_slides_by_phylum(self, image_guide_presentation):
        records = [
            {'observation_uuid': 1, 'scientific_name': 'A', 'phylum': 'Chordata'},
            {'observation_uuid': 2, 'scientific_name': 'B', 'phylum': 'Chordata'},
            {'observation_uuid': 3, 'scientific_name': 'C', 'phylum': 'Mollusca'},
        ]

        with patch.object(
                ImageGuidePresentation, '_fetch_normalized_image',
                side_effect=lambda localization: make_jpeg_bytes(),
        ):
            presentation = image_guide_presentation.build(records)

        assert len(presentation.slides) == 2
        chordata_slide, mollusca_slide = presentation.slides
        assert text_boxes_with_text(chordata_slide, 'C H O R D A T A')
        assert len(pictures(chordata_slide)) == 2
        assert text_boxes_with_text(mollusca_slide, 'M O L L U S C A')
        assert len(pictures(mollusca_slide)) == 1

    def test_build_starts_a_new_slide_after_six_records(self, image_guide_presentation):
        records = [
            {'observation_uuid': i, 'scientific_name': 'A', 'phylum': 'Chordata'} for i in range(7)
        ]

        with patch.object(
                ImageGuidePresentation, '_fetch_normalized_image',
                side_effect=lambda localization: make_jpeg_bytes(),
        ):
            presentation = image_guide_presentation.build(records)

        assert len(presentation.slides) == 2
        first_slide, second_slide = presentation.slides
        assert len(pictures(first_slide)) == 6
        assert len(pictures(second_slide)) == 1
        assert text_boxes_with_text(first_slide, 'C H O R D A T A')
        assert text_boxes_with_text(second_slide, 'C H O R D A T A')

    def test_build_skips_records_whose_image_fetch_fails(self, image_guide_presentation):
        records = [
            {'observation_uuid': 1, 'scientific_name': 'A', 'phylum': 'Chordata'},
            {'observation_uuid': 2, 'scientific_name': 'B', 'phylum': 'Chordata'},
            {'observation_uuid': 3, 'scientific_name': 'C', 'phylum': 'Chordata'},
        ]

        def fetch_or_fail(localization):
            if localization['observation_uuid'] == 2:
                raise ValueError('oh no!')
            return make_jpeg_bytes()

        with patch.object(
                ImageGuidePresentation, '_fetch_normalized_image', side_effect=fetch_or_fail,
        ):
            presentation = image_guide_presentation.build(records)

        assert len(presentation.slides) == 1
        assert len(pictures(presentation.slides[0])) == 2

    def test_fetch_all_images_keys_by_observation_uuid_despite_out_of_order_completion(
            self, image_guide_presentation,
    ):
        records = [{'observation_uuid': f'uuid-{i}'} for i in range(5)]

        # delays are deliberately inverted so results complete in the opposite order they were submitted in
        def fetch(localization):
            index = int(localization['observation_uuid'].split('-')[1])
            time.sleep((len(records) - index) * 0.01)
            return f'image-{localization["observation_uuid"]}'

        with patch.object(ImageGuidePresentation, '_fetch_normalized_image', side_effect=fetch):
            images = image_guide_presentation._fetch_all_images(records)

        assert images == {f'uuid-{i}': f'image-uuid-{i}' for i in range(5)}

    def test_fetch_all_images_omits_entry_for_failed_fetch_without_disturbing_others(
            self, image_guide_presentation,
    ):
        records = [{'observation_uuid': f'uuid-{i}'} for i in range(3)]

        def fetch(localization):
            if localization['observation_uuid'] == 'uuid-1':
                raise ValueError('oh no!')
            return f'image-{localization["observation_uuid"]}'

        with patch.object(ImageGuidePresentation, '_fetch_normalized_image', side_effect=fetch):
            images = image_guide_presentation._fetch_all_images(records)

        assert images == {'uuid-0': 'image-uuid-0', 'uuid-2': 'image-uuid-2'}

    @pytest.mark.parametrize('localization_type,attracted,expect_overlay', [
        (TatorLocalizationType.DOT, None, True),  # dropcam, no attracted value at all -> overlay
        (TatorLocalizationType.DOT, 'Not Attracted', True),  # dropcam, explicitly not attracted -> overlay
        (TatorLocalizationType.DOT, 'Attracted', False),  # dropcam, attracted -> no overlay
        (TatorLocalizationType.SUB_DOT, None, False),  # not a dropcam type -> no overlay regardless
    ])
    def test_build_adds_not_attracted_overlay_only_for_dropcam_non_attracted_records(
            self, image_guide_presentation, localization_type, attracted, expect_overlay,
    ):
        records = [
            {
                'observation_uuid': 1,
                'scientific_name': 'A',
                'phylum': 'Chordata',
                'type': localization_type,
                'attracted': attracted,
            },
        ]

        with patch.object(
                ImageGuidePresentation, '_fetch_normalized_image',
                side_effect=lambda localization: make_jpeg_bytes(),
        ):
            presentation = image_guide_presentation.build(records)

        overlay_present = bool(text_boxes_with_text(presentation.slides[0], 'NOT ATTRACTED'))
        assert overlay_present == expect_overlay

    def test_build_uses_unknown_phylum_fallback(self, image_guide_presentation):
        records = [{'observation_uuid': 1, 'scientific_name': 'A'}]  # no 'phylum' key

        with patch.object(
                ImageGuidePresentation, '_fetch_normalized_image',
                side_effect=lambda localization: make_jpeg_bytes(),
        ):
            presentation = image_guide_presentation.build(records)

        assert text_boxes_with_text(presentation.slides[0], 'U N K N O W N   P H Y L U M')

    def test_add_phylum_header_uppercases_and_spaces_out_letters(self, image_guide_presentation):
        slide = make_slide()

        image_guide_presentation._add_phylum_header(slide, 'Chordata')

        text_box = slide.shapes[0]
        assert text_box.text_frame.text == 'C H O R D A T A'

    def test_add_not_attracted_overlay_adds_red_bold_text(self, image_guide_presentation):
        slide = make_slide()

        image_guide_presentation._add_not_attracted_overlay(slide, Inches(0.5), Inches(1.5))

        text_box = slide.shapes[0]
        run = text_box.text_frame.paragraphs[0].runs[0]
        assert run.text == 'NOT ATTRACTED'
        assert run.font.bold is True
        assert str(run.font.color.rgb) == 'FF0000'

    def test_add_image_adds_picture_with_border(self, image_guide_presentation):
        slide = make_slide()

        image_guide_presentation._add_image(slide, make_jpeg_bytes(), Inches(0.5), Inches(1.5))

        picture = slide.shapes[0]
        assert picture.width == ImageGuidePresentation.IMAGE_WIDTH
        assert picture.height == ImageGuidePresentation.IMAGE_HEIGHT
        assert str(picture.line.color.rgb) == '000000'

    def test_make_run_sets_font_and_italic(self):
        slide = make_slide()
        text_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(1), Inches(1))
        paragraph = text_box.text_frame.paragraphs[0]

        run = ImageGuidePresentation._make_run(paragraph, 'Squalus', italic=True)

        assert run.text == 'Squalus'
        assert run.font.italic is True
        assert run.font.name == 'Arial'

    def test_add_image_header_plain_scientific_name(self, image_guide_presentation):
        slide = make_slide()
        localization = {'scientific_name': 'Cydippida'}

        image_guide_presentation._add_image_header(slide, localization, Inches(0.5), Inches(1.5))

        paragraph = slide.shapes[1].text_frame.paragraphs[0]
        assert run_texts(paragraph) == [('Cydippida', False)]

    def test_add_image_header_genus_without_species_appends_sp(self, image_guide_presentation):
        slide = make_slide()
        localization = {'scientific_name': 'Squalus', 'genus': 'Squalus'}

        image_guide_presentation._add_image_header(slide, localization, Inches(0.5), Inches(1.5))

        paragraph = slide.shapes[1].text_frame.paragraphs[0]
        assert run_texts(paragraph) == [('Squalus', True), (' sp.', False)]

    def test_add_image_header_genus_and_species_no_sp(self, image_guide_presentation):
        slide = make_slide()
        localization = {'scientific_name': 'Squalus acanthias', 'genus': 'Squalus', 'species': 'Squalus acanthias'}

        image_guide_presentation._add_image_header(slide, localization, Inches(0.5), Inches(1.5))

        paragraph = slide.shapes[1].text_frame.paragraphs[0]
        assert run_texts(paragraph) == [('Squalus acanthias', True)]

    def test_add_image_header_tentative_id_with_family(self, image_guide_presentation):
        slide = make_slide()
        localization = {
            'scientific_name': 'Cydippida',
            'tentative_id': 'Bathocyroe fosteri',
            'family': 'Bathocyroidae',
        }

        image_guide_presentation._add_image_header(slide, localization, Inches(0.5), Inches(1.5))

        paragraph = slide.shapes[1].text_frame.paragraphs[0]
        assert run_texts(paragraph) == [
            ('Cydippida', False),
            (' (', False),
            ('Bathocyroe fosteri', True),
            (' sp.', False),
            ('?)', False),
        ]

    def test_add_image_header_morphospecies_no_tentative_id(self, image_guide_presentation):
        slide = make_slide()
        localization = {'scientific_name': 'Cydippida', 'morphospecies': 'sp. 1', 'family': 'Bathocyroidae'}

        image_guide_presentation._add_image_header(slide, localization, Inches(0.5), Inches(1.5))

        paragraph = slide.shapes[1].text_frame.paragraphs[0]
        # morphospecies present -> no ' sp.' appended, and closing paren has no '?'
        assert run_texts(paragraph) == [
            ('Cydippida', False),
            (' (', False),
            ('sp. 1', True),
            (')', False),
        ]

    def test_fetch_normalized_image_crops_without_expansion_when_box_matches_target_aspect(
            self, image_guide_presentation,
    ):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.25, 0.25], 'dimensions': [0.5, 0.5]}],
        }
        source_image = make_jpeg_bytes(width=1600, height=900, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        assert Image.open(result).size == (800, 450)

    def test_fetch_normalized_image_expands_narrow_box_horizontally(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.4, 0.4], 'dimensions': [0.2, 0.3]}],
        }
        source_image = make_jpeg_bytes(width=1000, height=1000, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # box was 200x300 (aspect 0.67); widened to match 16:9 while keeping the original height
        assert Image.open(result).size == (533, 300)

    def test_fetch_normalized_image_expands_wide_box_vertically(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.3, 0.4], 'dimensions': [0.4, 0.1]}],
        }
        source_image = make_jpeg_bytes(width=1000, height=1000, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # box was 400x100 (aspect 4.0); heightened to match 16:9 while keeping the original width
        assert Image.open(result).size == (400, 225)

    def test_fetch_normalized_image_clamps_crop_at_left_edge(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.05, 0.4], 'dimensions': [0.1, 0.3]}],
        }
        source_image = make_jpeg_bytes(width=500, height=500, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # expansion would push left of frame 0; clamped to 0 and shifted right to preserve width
        assert Image.open(result).size == (266, 150)

    def test_fetch_normalized_image_clamps_crop_at_right_edge(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.85, 0.4], 'dimensions': [0.1, 0.3]}],
        }
        source_image = make_jpeg_bytes(width=500, height=500, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # expansion would push past the right edge of the frame; clamped and shifted left to preserve width
        assert Image.open(result).size == (267, 150)

    def test_fetch_normalized_image_clamps_crop_at_top_edge(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.3, 0.02], 'dimensions': [0.4, 0.05]}],
        }
        source_image = make_jpeg_bytes(width=500, height=500, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # expansion would push above the top edge of the frame; clamped to 0 and shifted down to preserve height
        assert Image.open(result).size == (200, 111)

    def test_fetch_normalized_image_clamps_crop_at_bottom_edge(self, image_guide_presentation):
        localization = {
            'media_id': 1,
            'frame': 1,
            'all_localizations': [{'points': [0.3, 0.93], 'dimensions': [0.4, 0.05]}],
        }
        source_image = make_jpeg_bytes(width=500, height=500, color='blue')

        with patch.object(TatorRestClient, 'get_frame', return_value=source_image.getvalue()):
            result = image_guide_presentation._fetch_normalized_image(localization)

        # expansion would push past the bottom edge of the frame; clamped and shifted up to preserve height
        assert Image.open(result).size == (200, 112)
